from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from .lesson_pages import LessonDocument

_SLIDE_WIDTH_IN = 7.5
_SLIDE_HEIGHT_IN = 10.0


def write_pptx_export(path: str | Path, document: LessonDocument, rendered_image_paths: list[Path]) -> None:
    """1ページ=1スライドのPPTXを書き出す（完成outputとしてのPPTX export）。

    スライドのレイアウト編集そのものはPowerPoint/Canva等で行う前提であり、ここでは
    「各ページの完成画像をスライドに配置し、タイトルを添える」という最小限の構成に留める。
    スライド内の複雑な図形・アニメーション等の再現は対象外（`rendered/`の画像自体が
    最終的な見た目を担う）。
    """
    output_path = Path(path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    presentation.slide_width = Inches(_SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(_SLIDE_HEIGHT_IN)
    blank_layout = presentation.slide_layouts[6]

    image_by_page_no = {
        int(image_path.stem.split("_")[-1]): image_path
        for image_path in rendered_image_paths
    }

    for page in document.pages:
        slide = presentation.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(_SLIDE_WIDTH_IN - 0.6), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"Page {page.page_no}: {page.title}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True

        image_path = image_by_page_no.get(page.page_no)
        if image_path and image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(0.5),
                Inches(1.0),
                width=Inches(_SLIDE_WIDTH_IN - 1.0),
            )
        else:
            body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(_SLIDE_WIDTH_IN - 1.0), Inches(_SLIDE_HEIGHT_IN - 1.5))
            body_box.text_frame.text = page.summary or page.title

    presentation.save(output_path)
