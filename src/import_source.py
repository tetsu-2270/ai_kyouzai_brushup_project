from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


def _try_ocr(image_path: Path) -> str:
    """画像からテキストを抽出する。pytesseractまたはtesseract本体が無い環境では空文字を返す。

    OCR精度の向上は対象外。まずは画像を入力として扱えること・画像を落とさないことを優先し、
    OCRが使えない環境でも取り込み自体は成立するようにする。
    """
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        return ""
    try:
        with Image.open(image_path) as image:
            return pytesseract.image_to_string(image, lang="jpn+eng")
    except Exception:
        return ""


def _derive_title_and_summary(text: str, index: int, fallback_name: str) -> tuple[str, str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if lines:
        title = lines[0][:60]
        summary = " ".join(lines[:2])[:120]
    else:
        title = f"取り込みページ {index}（{fallback_name}）"
        summary = ""
    return title, summary


def _text_to_lines(text: str) -> list[dict[str, str]]:
    return [{"speaker": "", "text": line.strip()} for line in text.splitlines() if line.strip()]


def _copy_asset(src: Path, assets_dir: Path, dest_name: str) -> str:
    assets_dir.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(src, assets_dir / dest_name)
    return f"assets/{dest_name}"


def _page_from_image(index: int, image_path: Path, assets_dir: Path) -> dict[str, Any]:
    ocr_text = _try_ocr(image_path)
    title, summary = _derive_title_and_summary(ocr_text, index, image_path.name)
    dest_name = f"page_{index:03d}{image_path.suffix.lower()}"
    source_image = _copy_asset(image_path, assets_dir, dest_name)
    main_visual = f"元画像（{source_image}）を参照してレイアウトを作成する。"
    notes = "" if ocr_text.strip() else "OCRでテキストを抽出できませんでした。元画像を直接参照してください。"
    return {
        "page_no": index,
        "source_image": source_image,
        "source_assets": [],
        "title": title,
        "summary": summary,
        "lines": _text_to_lines(ocr_text),
        "improvement_points": [],
        "canva": {"layout_type": "", "main_visual": main_visual, "notes": notes},
    }


def import_images(image_paths: list[Path], assets_dir: Path, project_title: str) -> dict[str, Any]:
    """画像ファイル群を、ファイル名順に1画像=1元ページとして取り込む。"""
    sorted_paths = sorted(image_paths, key=lambda p: p.name)
    pages = [
        _page_from_image(index, path, assets_dir)
        for index, path in enumerate(sorted_paths, start=1)
    ]
    return {"project_title": project_title, "target_reader": "教材制作者", "pages": pages}


def import_pdf(pdf_path: Path, assets_dir: Path) -> dict[str, Any]:
    """PDFをページ単位で取り込む。各ページをテキスト抽出しつつ、ページ画像も保存する（PyMuPDF使用）。"""
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise ValueError(
            "PDF取り込みにはPyMuPDF(pymupdf)が必要です。`python3 -m pip install pymupdf`でインストールしてください。"
        ) from e

    pages: list[dict[str, Any]] = []
    assets_dir.mkdir(parents=True, exist_ok=True)
    with fitz.open(pdf_path) as doc:
        for index, page in enumerate(doc, start=1):
            dest_name = f"page_{index:03d}.png"
            page.get_pixmap().save(assets_dir / dest_name)
            source_image = f"assets/{dest_name}"

            text = page.get_text().strip()
            title, summary = _derive_title_and_summary(text, index, pdf_path.name)
            main_visual = f"元ページ画像（{source_image}）を参照してレイアウトを作成する。"
            pages.append({
                "page_no": index,
                "source_image": source_image,
                "source_assets": [],
                "title": title,
                "summary": summary,
                "lines": _text_to_lines(text),
                "improvement_points": [],
                "canva": {"layout_type": "", "main_visual": main_visual, "notes": ""},
            })
    return {"project_title": pdf_path.stem, "target_reader": "教材制作者", "pages": pages}


_PPTX_SLIDE_NOTE = "スライド全体の画像化（レンダリング）は未対応です。スライド内に埋め込まれた画像のみ保持しています。"


def import_pptx(pptx_path: Path, assets_dir: Path) -> dict[str, Any]:
    """PPTXをスライド単位で取り込む。スライド内テキストと埋め込み画像を抽出する。

    スライド全体を1枚のビジュアルとしてレンダリングするには外部レンダラー（PowerPoint/LibreOffice等）
    が必要であり、今回は対象外とする。代わりにスライド内に埋め込まれた画像をsource_image/
    source_assetsとして保持し、「画像を落とさない」方針を満たす。
    """
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ValueError(
            "PPTX取り込みにはpython-pptxが必要です。`python3 -m pip install python-pptx`でインストールしてください。"
        ) from e

    presentation = Presentation(pptx_path)
    pages: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        asset_paths: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            image = getattr(shape, "image", None)
            if image is not None:
                assets_dir.mkdir(parents=True, exist_ok=True)
                dest_name = f"slide_{index:03d}_{len(asset_paths) + 1}.{image.ext}"
                (assets_dir / dest_name).write_bytes(image.blob)
                asset_paths.append(f"assets/{dest_name}")

        slide_text = "\n".join(texts)
        title, summary = _derive_title_and_summary(slide_text, index, pptx_path.name)
        source_image = asset_paths[0] if asset_paths else ""
        source_assets = asset_paths[1:]
        if source_image:
            main_visual = f"スライド内の画像（{source_image}）を参照してレイアウトを作成する。"
        else:
            main_visual = "このスライドには埋め込み画像がありません。テキストのみを参考にレイアウトを作成する。"
        pages.append({
            "page_no": index,
            "source_image": source_image,
            "source_assets": source_assets,
            "title": title,
            "summary": summary,
            "lines": _text_to_lines(slide_text),
            "improvement_points": [],
            "canva": {"layout_type": "", "main_visual": main_visual, "notes": _PPTX_SLIDE_NOTE},
        })
    return {"project_title": pptx_path.stem, "target_reader": "教材制作者", "pages": pages}


def import_source(input_path: str | Path, assets_dir: str | Path) -> dict[str, Any]:
    """元資料（画像ディレクトリ/画像ファイル/PDF/PPTX）を、pages形式JSON互換の辞書として取り込む。

    戻り値はexamples/sample_pages.jsonと同じpages形式であり、そのままlesson-pagesの
    --inputに渡せる。画像アセットはassets_dir配下にコピー・保存する。
    """
    path = Path(input_path)
    assets_dir = Path(assets_dir)
    if not path.exists():
        raise FileNotFoundError(f"取り込み元が見つかりません: {path}")

    if path.is_dir():
        image_paths = [p for p in path.iterdir() if p.is_file() and p.suffix.lower() in _IMAGE_EXTENSIONS]
        if not image_paths:
            raise ValueError(
                f"{path} 配下に画像ファイル（.png/.jpg/.jpeg/.webp）が見つかりません。"
            )
        return import_images(image_paths, assets_dir, project_title=path.name)

    suffix = path.suffix.lower()
    if suffix in _IMAGE_EXTENSIONS:
        return import_images([path], assets_dir, project_title=path.stem)
    if suffix == ".pdf":
        return import_pdf(path, assets_dir)
    if suffix == ".pptx":
        return import_pptx(path, assets_dir)
    if suffix == ".ppt":
        raise ValueError(
            ".ppt（旧形式）は現時点では未対応です。PowerPoint等で.pptxに変換してから再度お試しください。"
        )
    raise ValueError(
        f"対応していない形式です: {suffix}。対応形式は画像（.png/.jpg/.jpeg/.webp）/.pdf/.pptxです。"
    )
