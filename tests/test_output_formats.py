import json
from pathlib import Path

import pytest
from PIL import Image

from src.cli import main


def _make_source_images(source_dir, count=2):
    source_dir.mkdir(parents=True, exist_ok=True)
    for i in range(1, count + 1):
        Image.new("RGB", (80, 120), color=(220, 220, 220)).save(source_dir / f"page_{i:03d}.png")


def _make_source_pdf(pdf_path, page_count=2):
    import fitz

    doc = fitz.open()
    for i in range(page_count):
        page = doc.new_page()
        page.insert_text((72, 72), f"page {i + 1}", fontsize=14)
    doc.save(pdf_path)
    doc.close()


def _make_source_pptx(pptx_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "スライドタイトル"
    prs.save(pptx_path)


def test_build_all_always_writes_editable_lesson_pages_json(tmp_path, monkeypatch):
    """--output-formatの指定に関わらず、editable/lesson_pages.jsonは常に生成される。"""
    source_dir = tmp_path / "source"
    _make_source_images(source_dir)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    editable_path = output_dir / "editable" / "lesson_pages.json"
    assert editable_path.exists()
    data = json.loads(editable_path.read_text(encoding="utf-8"))
    assert len(data["pages"]) == 2
    # Phase 8互換のlesson_pages.json/canva_design.mdはoutput/compat/配下に生成される
    # （editable//canva/と同名のファイルをoutput_dir直下に重複させないため）。
    assert (output_dir / "compat" / "lesson_pages.json").exists()
    assert (output_dir / "compat" / "canva_design.md").exists()


def test_build_all_does_not_duplicate_lesson_pages_json_or_canva_design_md_at_root(tmp_path, monkeypatch):
    """output/lesson_pages.jsonとoutput/canva_design.mdが直下に重複生成されないことを確認する。

    正式な編集対象はoutput/editable/lesson_pages.jsonのみ、正式なCanva指示書は
    output/canva/canva_design.mdのみであり、output_dir直下には同名ファイルを置かない。
    """
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "all",
        ],
    )
    main()

    assert not (output_dir / "lesson_pages.json").exists()
    assert not (output_dir / "canva_design.md").exists()
    assert (output_dir / "editable" / "lesson_pages.json").exists()
    assert (output_dir / "canva" / "canva_design.md").exists()
    assert (output_dir / "compat" / "lesson_pages.json").exists()
    assert (output_dir / "compat" / "canva_design.md").exists()


def test_build_all_no_compat_output_flag_skips_compat_directory(tmp_path, monkeypatch):
    """--no-compat-outputを指定すると、output/compat/配下(lesson_pages.json/canva_design.md/
    brushup.md/brushup.docx/brushup.pdf)が一切生成されないことを確認する。

    editable//canva/の生成には影響しない。
    """
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--no-compat-output",
        ],
    )
    main()

    assert not (output_dir / "compat").exists()
    assert not (output_dir / "brushup.md").exists()
    assert not (output_dir / "brushup.docx").exists()
    assert not (output_dir / "brushup.pdf").exists()
    assert (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_does_not_duplicate_brushup_role_at_root(tmp_path, monkeypatch):
    """brushup.*(md/docx/pdf)がoutput_dir直下に生成されず、正式outputはexports/、
    後方互換outputはcompat/に整理されていることを確認する（Phase 9.2）。
    """
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "all",
        ],
    )
    main()

    # output_dir直下にはbrushup.*を置かない。
    assert not (output_dir / "brushup.md").exists()
    assert not (output_dir / "brushup.docx").exists()
    assert not (output_dir / "brushup.pdf").exists()

    # 正式な完成outputはexports/にある。
    assert (output_dir / "exports" / "material.md").exists()
    assert (output_dir / "exports" / "material.docx").exists()
    assert (output_dir / "exports" / "material.pdf").exists()

    # 後方互換outputはcompat/にまとめられる（既定で生成）。
    assert (output_dir / "compat" / "brushup.md").exists()
    assert (output_dir / "compat" / "brushup.docx").exists()
    assert (output_dir / "compat" / "brushup.pdf").exists()

    # Phase 9.1で整理済みのcompat/lesson_pages.json・canva_design.mdの扱いは変わらない。
    assert (output_dir / "compat" / "lesson_pages.json").exists()
    assert (output_dir / "compat" / "canva_design.md").exists()
    assert (output_dir / "editable" / "lesson_pages.json").exists()
    assert (output_dir / "canva" / "canva_design.md").exists()


def test_build_all_no_compat_output_skips_brushup_compat_files(tmp_path, monkeypatch):
    """--no-compat-output指定時は、compat/brushup.*が一切生成されないことを確認する。"""
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "all",
            "--no-compat-output",
        ],
    )
    main()

    assert not (output_dir / "compat").exists()
    assert not (output_dir / "brushup.md").exists()
    assert not (output_dir / "brushup.docx").exists()
    assert not (output_dir / "brushup.pdf").exists()
    # 正式outputは--no-compat-outputの影響を受けない。
    assert (output_dir / "exports" / "material.md").exists()
    assert (output_dir / "editable" / "lesson_pages.json").exists()
    assert (output_dir / "canva" / "canva_design.md").exists()


def test_build_all_default_output_format_resolves_to_image_for_image_input(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    assert (output_dir / "rendered" / "page_001.png").exists()
    assert (output_dir / "rendered" / "page_002.png").exists()
    # 明示的にcanva/pptx/docx/mdを要求していないため、それらは生成されない。
    assert not (output_dir / "canva").exists()
    assert not (output_dir / "exports").exists()


def test_build_all_default_output_format_resolves_to_pdf_for_pdf_input(tmp_path, monkeypatch):
    pdf_path = tmp_path / "source.pdf"
    _make_source_pdf(pdf_path)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(pdf_path), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    assert (output_dir / "exports" / "material.pdf").exists()
    assert not (output_dir / "rendered").exists()


def test_build_all_default_output_format_resolves_to_pptx_for_pptx_input(tmp_path, monkeypatch):
    pptx_path = tmp_path / "source.pptx"
    _make_source_pptx(pptx_path)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(pptx_path), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    assert (output_dir / "exports" / "material.pptx").exists()


def test_build_all_output_format_image_generates_rendered_pages(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=3)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "image",
        ],
    )
    main()

    for i in (1, 2, 3):
        assert (output_dir / "rendered" / f"page_{i:03d}.png").exists()


def test_build_all_output_format_canva_generates_canva_option_output(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "canva",
        ],
    )
    main()

    canva_path = output_dir / "canva" / "canva_design.md"
    assert canva_path.exists()
    assert "元画像: assets/page_001.png" in canva_path.read_text(encoding="utf-8")
    # canva指定時でも、後方互換のためcompat/配下の一式は引き続き生成される
    # （Canva指示書だけが最終成果物になる設計にはしない）。output_dir直下には生成しない。
    assert (output_dir / "compat" / "brushup.md").exists()
    assert (output_dir / "compat" / "brushup.docx").exists()
    assert not (output_dir / "brushup.md").exists()
    assert not (output_dir / "brushup.docx").exists()


def test_build_all_output_format_json_does_not_create_rendered_or_exports(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "json",
        ],
    )
    main()

    assert (output_dir / "editable" / "lesson_pages.json").exists()
    assert not (output_dir / "rendered").exists()
    assert not (output_dir / "exports").exists()
    assert not (output_dir / "canva").exists()


def test_build_all_output_format_all_generates_every_format(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "all",
        ],
    )
    main()

    assert (output_dir / "editable" / "lesson_pages.json").exists()
    assert (output_dir / "rendered" / "page_001.png").exists()
    assert (output_dir / "canva" / "canva_design.md").exists()
    assert (output_dir / "exports" / "material.pdf").exists()
    assert (output_dir / "exports" / "material.pptx").exists()
    assert (output_dir / "exports" / "material.docx").exists()
    assert (output_dir / "exports" / "material.md").exists()


def test_build_all_output_format_pdf_pptx_docx_md_individually(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)

    for output_format, expected_file in (
        ("pdf", "material.pdf"),
        ("pptx", "material.pptx"),
        ("docx", "material.docx"),
        ("md", "material.md"),
    ):
        output_dir = tmp_path / f"output_{output_format}"
        monkeypatch.setattr(
            "sys.argv",
            [
                "cli", "build-all",
                "--input", str(source_dir),
                "--mode", "proofread",
                "--output-dir", str(output_dir),
                "--output-format", output_format,
            ],
        )
        main()
        assert (output_dir / "exports" / expected_file).exists()


def test_regenerate_from_editable_file_produces_rendered_images(tmp_path, monkeypatch):
    """output/editable/lesson_pages.json を編集した後、再生成コマンドでrendered/を作り直せる。"""
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    editable_path = output_dir / "editable" / "lesson_pages.json"
    data = json.loads(editable_path.read_text(encoding="utf-8"))
    data["pages"][0]["title"] = "編集後のタイトル"
    editable_path.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "regenerate", "--input", str(editable_path), "--output-format", "canva"],
    )
    main()

    canva_text = (output_dir / "canva" / "canva_design.md").read_text(encoding="utf-8")
    assert "編集後のタイトル" in canva_text


def test_regenerate_with_explicit_output_dir(tmp_path, monkeypatch):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    editable_path = output_dir / "editable" / "lesson_pages.json"
    other_output_dir = tmp_path / "other_output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "regenerate",
            "--input", str(editable_path),
            "--output-format", "image",
            "--output-dir", str(other_output_dir),
        ],
    )
    main()

    assert (other_output_dir / "rendered" / "page_001.png").exists()


def test_regenerate_works_for_generate_mode_document_without_source_image(tmp_path, monkeypatch):
    """新規構築(generateモード)相当のsource_image無しドキュメントでも、再生成が例外なく成功する。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages",
            "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "regenerate", "--input", str(editable_path), "--output-format", "all"],
    )
    main()

    output_dir = editable_path.parent.parent
    assert (output_dir / "rendered").exists()
    assert len(list((output_dir / "rendered").glob("page_*.png"))) > 0
    assert (output_dir / "canva" / "canva_design.md").exists()
    canva_text = (output_dir / "canva" / "canva_design.md").read_text(encoding="utf-8")
    assert "元画像:" not in canva_text
    assert "参考画像:" not in canva_text


# --- Phase 10: --font-path / フォント未検出警告 ----------------------------------


def test_build_all_accepts_font_path_option(tmp_path, monkeypatch):
    """build-all --font-path が受け付けられ、指定フォントで画像outputが生成されることを確認する。"""
    from src.image_renderer import resolve_font_path

    font_path = resolve_font_path(None)
    if font_path is None:
        return  # 実行環境に日本語フォントが無い場合はスキップ相当

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "image",
            "--font-path", font_path,
        ],
    )
    main()

    assert (output_dir / "rendered" / "page_001.png").exists()


def test_build_all_font_path_reaches_image_renderer(tmp_path, monkeypatch):
    """--font-pathで指定したパスが実際にimage_rendererへ渡ることを確認する（呼び出し引数を検証）。"""
    import src.cli as cli_module

    captured = {}
    original = cli_module.render_document_images

    def _spy(document, output_dir, rendered_dir, font_path=None):
        captured["font_path"] = font_path
        return original(document, output_dir, rendered_dir, font_path=font_path)

    monkeypatch.setattr(cli_module, "render_document_images", _spy)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "image",
            "--font-path", "/dummy/font/path.ttc",
        ],
    )
    try:
        main()
    except SystemExit:
        pass  # 存在しないダミーパスなのでエラー終了するが、渡された引数の検証が目的
    assert captured.get("font_path") == "/dummy/font/path.ttc"


def test_build_all_invalid_font_path_raises_clear_error(tmp_path, monkeypatch, capsys):
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all",
            "--input", str(source_dir),
            "--mode", "proofread",
            "--output-dir", str(output_dir),
            "--output-format", "image",
            "--font-path", "/no/such/font.ttc",
        ],
    )
    try:
        main()
        assert False, "SystemExitが発生するはず"
    except SystemExit as e:
        assert e.code == 1
    captured = capsys.readouterr()
    assert "見つかりません" in captured.err


def test_regenerate_accepts_font_path_option(tmp_path, monkeypatch):
    """regenerate --font-path が受け付けられ、指定フォントで画像outputが再生成されることを確認する。"""
    from src.image_renderer import resolve_font_path

    font_path = resolve_font_path(None)
    if font_path is None:
        return  # 実行環境に日本語フォントが無い場合はスキップ相当

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages",
            "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "regenerate",
            "--input", str(editable_path),
            "--output-format", "image",
            "--font-path", font_path,
        ],
    )
    main()

    output_dir = editable_path.parent.parent
    assert len(list((output_dir / "rendered").glob("page_*.png"))) > 0


def test_build_all_warns_when_synthesizing_text_without_japanese_font(tmp_path, monkeypatch, capsys):
    """source_imageが無いページの画像合成が必要なのに日本語フォントが無い場合、警告が出ることを確認する。"""
    import src.image_renderer as image_renderer_module

    monkeypatch.setattr(image_renderer_module, "_JAPANESE_FONT_CANDIDATES", ())

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages",
            "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "regenerate", "--input", str(editable_path), "--output-format", "image"],
    )
    main()

    captured = capsys.readouterr()
    assert "WARNING" in captured.err
    # 警告が出ても画像自体は生成され、処理が継続していることを確認する。
    output_dir = editable_path.parent.parent
    assert len(list((output_dir / "rendered").glob("page_*.png"))) > 0


# --- Phase 10.1: OCR前提の事前チェック・結果反映 ----------------------------------


def test_cli_check_ocr_runs_and_prints_report(capsys):
    """check-ocrコマンドが例外なく実行され、診断レポートを出力することを確認する。"""
    monkeypatch_argv = ["cli", "check-ocr"]
    import sys as sys_module

    original_argv = sys_module.argv
    sys_module.argv = monkeypatch_argv
    try:
        main()
    finally:
        sys_module.argv = original_argv

    captured = capsys.readouterr()
    assert "OCR" in captured.out


def test_build_all_fails_when_tesseract_missing_for_image_input_proofread(tmp_path, monkeypatch, capsys):
    """画像input + proofreadモードでtesseractが無い場合、build-allがエラー終了することを確認する
    （Phase 10.1追加修正: OCR必須モードでOCR不能なまま空データで成功させない）。"""
    import src.cli as cli_module
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: not_ready_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()

    assert exc_info.value.code == 1
    captured = capsys.readouterr()
    assert "ERROR" in captured.err
    assert "Tesseract" in captured.err
    assert "proofread" in captured.err
    # editable/lesson_pages.jsonはOCR前提チェックより後の工程で生成されるため、作られない。
    assert not (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_fails_when_japanese_language_data_missing_for_proofread(tmp_path, monkeypatch, capsys):
    """画像input + proofreadモードでjpn言語データが無い場合、build-allがエラー終了することを確認する。"""
    import src.cli as cli_module
    import src.import_source as import_source_module

    jpn_missing_status = {
        "tesseract_available": True, "tesseract_path": "/usr/bin/tesseract", "tesseract_on_path": True,
        "languages": ["eng"], "japanese_available": False, "english_available": True,
        "brew_available": True, "brew_path": "/usr/local/bin/brew", "brew_on_path": True,
        "path_suggestions": [], "warnings": [], "errors": ["Japanese OCR language data 'jpn' was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: jpn_missing_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: jpn_missing_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()

    assert exc_info.value.code == 1
    captured = capsys.readouterr()
    assert "ERROR" in captured.err
    assert "jpn" in captured.err
    assert not (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_fails_when_all_pages_ocr_empty_for_proofread(tmp_path, monkeypatch, capsys):
    """Tesseract自体は使えるが、全ページのOCR結果が空の場合もbuild-allがエラー終了することを確認する。"""
    import src.import_source as import_source_module

    monkeypatch.setattr(import_source_module, "_try_ocr", lambda image_path, ocr_status: "")

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()

    assert exc_info.value.code == 1
    captured = capsys.readouterr()
    assert "ERROR" in captured.err
    assert "proofread" in captured.err
    assert not (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_warns_and_continues_when_some_pages_ocr_empty(tmp_path, monkeypatch, capsys):
    """一部のページだけOCR結果が空の場合は、警告のうえ処理を継続することを確認する。"""
    import src.import_source as import_source_module

    call_count = {"n": 0}

    def _fake_ocr(image_path, ocr_status):
        call_count["n"] += 1
        return "検出されたテキスト" if call_count["n"] == 1 else ""

    monkeypatch.setattr(import_source_module, "_try_ocr", _fake_ocr)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    captured = capsys.readouterr()
    assert "WARNING" in captured.err
    assert "1 of 2" in captured.err
    assert (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_allow_empty_ocr_flag_bypasses_ocr_precondition(tmp_path, monkeypatch, capsys):
    """--allow-empty-ocrを指定すれば、Tesseract未導入・全ページ空でもエラー終了しないことを確認する。"""
    import src.import_source as import_source_module

    monkeypatch.setattr(import_source_module, "_try_ocr", lambda image_path, ocr_status: "")

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all", "--input", str(source_dir), "--mode", "proofread",
            "--output-dir", str(output_dir), "--allow-empty-ocr",
        ],
    )
    main()

    assert (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_ocr_precondition_does_not_apply_to_pdf_input(tmp_path):
    """PDF inputはOCRではなくネイティブなテキスト抽出を使うため、OCR前提チェックの対象外であることを
    確認する（PDFにテキストがあれば、Tesseractの有無に関わらず処理が継続する）。"""
    import fitz

    pdf_path = tmp_path / "source.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "PDFのテキスト")
    doc.save(pdf_path)
    doc.close()

    output_dir = tmp_path / "output"
    import sys as sys_module

    original_argv = sys_module.argv
    sys_module.argv = [
        "cli", "build-all", "--input", str(pdf_path), "--mode", "proofread", "--output-dir", str(output_dir),
    ]
    try:
        main()
    finally:
        sys_module.argv = original_argv

    assert (output_dir / "editable" / "lesson_pages.json").exists()


def test_build_all_ocr_result_flows_into_editable_lesson_pages_json(tmp_path, monkeypatch):
    """OCRが成功した場合、その結果がimported_pages.json・editable/lesson_pages.jsonの
    bodyに反映され、proofreadの校正対象として使えることを確認する。"""
    import src.import_source as import_source_module

    monkeypatch.setattr(
        import_source_module, "_try_ocr",
        lambda image_path, ocr_status: "教材の重要なポイントです。",
    )

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    imported_data = json.loads((output_dir / "imported_pages.json").read_text(encoding="utf-8"))
    assert imported_data["pages"][0]["lines"] == [{"speaker": "", "text": "教材の重要なポイントです。"}]

    editable_data = json.loads((output_dir / "editable" / "lesson_pages.json").read_text(encoding="utf-8"))
    assert "教材の重要なポイントです。" in editable_data["pages"][0]["body"]


def test_build_all_does_not_warn_all_pages_empty_when_ocr_succeeds(tmp_path, monkeypatch, capsys):
    import src.import_source as import_source_module

    monkeypatch.setattr(
        import_source_module, "_try_ocr",
        lambda image_path, ocr_status: "抽出されたテキスト",
    )

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    captured = capsys.readouterr()
    assert "check-ocr" not in captured.err


# --- Phase 10.2: 実行ログ出力 -----------------------------------------------------


def _log_dir(monkeypatch) -> "Path":
    import os
    from pathlib import Path

    return Path(os.environ["AI_KYOUZAI_LOGS_DIR"])


def test_build_all_writes_log_file_on_success(tmp_path, monkeypatch):
    """build-all成功時、logs/にbuild-allのログファイルが作成されることを確認する。"""
    import fitz

    pdf_path = tmp_path / "source.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "sample text")
    doc.save(pdf_path)
    doc.close()

    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(pdf_path), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    log_files = list(_log_dir(monkeypatch).glob("*_build-all.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "command: build-all" in text
    assert "exit_code: 0" in text
    assert "INPUT" in text
    assert "OUTPUT" in text


def test_build_all_writes_log_file_on_failure(tmp_path, monkeypatch):
    """build-all失敗時（OCR不能）にも、exit_code・エラー内容を含むログが残ることを確認する。"""
    import src.cli as cli_module
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: not_ready_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit):
        main()

    log_files = list(_log_dir(monkeypatch).glob("*_build-all.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "exit_code: 1" in text
    assert "Tesseract" in text


def test_lesson_pages_generate_mode_writes_log_named_generate(tmp_path, monkeypatch):
    """lesson-pages --mode generateの実行ログが、モード名(generate)でファイル名になることを確認する。"""
    output_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(output_path),
        ],
    )
    main()

    log_files = list(_log_dir(monkeypatch).glob("*_generate.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "exit_code: 0" in text
    assert str(output_path) in text


def test_regenerate_writes_log_file_with_input_and_output_format(tmp_path, monkeypatch):
    """regenerate実行ログに、入力パスとoutput-formatが記録されることを確認する。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "regenerate", "--input", str(editable_path), "--output-format", "image"],
    )
    main()

    log_files = list(_log_dir(monkeypatch).glob("*_regenerate.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert str(editable_path) in text
    assert "image" in text
    assert "exit_code: 0" in text


def test_regenerate_writes_log_on_json_error(tmp_path, monkeypatch):
    """regenerateがJSON構文エラーで失敗した場合も、ログにエラー内容が残ることを確認する。"""
    broken_path = tmp_path / "broken.json"
    broken_path.write_text("{ this is not valid json", encoding="utf-8")

    monkeypatch.setattr(
        "sys.argv", ["cli", "regenerate", "--input", str(broken_path), "--output-format", "image"]
    )
    with pytest.raises(SystemExit):
        main()

    log_files = list(_log_dir(monkeypatch).glob("*_regenerate.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "exit_code: 1" in text
    assert "JSON" in text or "不正" in text


def test_check_ocr_writes_log_file_with_ocr_summary(tmp_path, monkeypatch):
    monkeypatch.setattr("sys.argv", ["cli", "check-ocr"])
    main()

    log_files = list(_log_dir(monkeypatch).glob("*_check-ocr.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "exit_code: 0" in text
    assert "tesseract_available" in text


# --- Phase 10.2: 成功判定の見直し ---------------------------------------------------


def test_build_all_fails_when_input_directory_empty(tmp_path, monkeypatch):
    """input/sourceが空ディレクトリの場合、build-allが非ゼロ終了することを確認する
    （import_source()の既存動作。回帰確認）。"""
    empty_source = tmp_path / "empty_source"
    empty_source.mkdir()
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(empty_source), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_build_all_fails_when_only_unsupported_files_in_directory(tmp_path, monkeypatch):
    """input/sourceに対応外ファイル（.txt等）しか無い場合、build-allが非ゼロ終了することを確認する。"""
    source_dir = tmp_path / "source"
    source_dir.mkdir()
    (source_dir / "notes.txt").write_text("これは画像でもPDFでもPPTXでもない", encoding="utf-8")
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_regenerate_fails_when_pages_empty(tmp_path, monkeypatch):
    """editable中間ファイルのpagesが空の場合、regenerateが非ゼロ終了することを確認する。"""
    empty_pages_path = tmp_path / "lesson_pages.json"
    empty_pages_path.write_text(
        json.dumps({
            "metadata": {
                "project_title": "空テスト", "mode": "generate", "source_policy": "generate",
                "target_audience": "", "tone": "", "generated_at": "2026-01-01T00:00:00",
                "requirements_source": None,
            },
            "pages": [],
        }, ensure_ascii=False),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "sys.argv", ["cli", "regenerate", "--input", str(empty_pages_path), "--output-format", "image"]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_regenerate_fails_when_specified_output_format_artifact_not_generated(tmp_path, monkeypatch):
    """レンダラーが何らかの理由で成果物を生成できなかった場合、regenerateが非ゼロ終了することを
    確認する（実質失敗を正常終了扱いにしないための安全網）。"""
    import src.cli as cli_module

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(cli_module, "write_pdf", lambda path, document: None)
    monkeypatch.setattr(
        "sys.argv", ["cli", "regenerate", "--input", str(editable_path), "--output-format", "pdf"]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


# --- Phase 10.2追加修正: 個別CLIの成果物未生成チェック -------------------------------


def test_lesson_pages_cli_fails_when_pages_empty(tmp_path, monkeypatch):
    """個別CLIのlesson-pagesが、pagesが空のpages形式JSONを渡された場合に非ゼロ終了することを
    確認する（build-all/regenerateと同様、実質失敗を正常終了扱いにしない）。"""
    empty_input = tmp_path / "empty_pages.json"
    empty_input.write_text(
        json.dumps({"project_title": "空テスト", "target_reader": "テスター", "pages": []}, ensure_ascii=False),
        encoding="utf-8",
    )
    output_path = tmp_path / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "lesson-pages", "--mode", "proofread", "--input", str(empty_input), "--output", str(output_path)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_import_source_cli_fails_when_imported_pages_are_empty(tmp_path, monkeypatch):
    """個別CLIのimport-sourceが、取り込み結果のpagesが0件だった場合に非ゼロ終了することを
    確認する（画像ディレクトリが空のケースは既存のimport_source()自身が検知するが、
    PPTXでスライドが1つも無い等のケースに備えた追加の安全網）。pymupdfはページ0件のPDFを
    保存できないため、import_source()自体をモックして「pagesが空の取り込み結果」を再現する。
    """
    import src.cli as cli_module

    monkeypatch.setattr(cli_module, "import_source", lambda input_path, assets_dir, quiet=False: {"pages": []})

    dummy_input = tmp_path / "dummy.pptx"
    dummy_input.write_bytes(b"dummy")
    output_path = tmp_path / "imported_pages.json"
    monkeypatch.setattr(
        "sys.argv", ["cli", "import-source", "--input", str(dummy_input), "--output", str(output_path)]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_canva_cli_fails_when_output_not_generated(tmp_path, monkeypatch):
    """個別CLIのcanvaが、レンダラーが空文字列を返す（＝実質何も生成できない）場合に
    非ゼロ終了することを確認する。"""
    import src.cli as cli_module

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(cli_module, "render_canva_design", lambda document: "")
    output_path = tmp_path / "canva_design.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "canva", "--input", str(editable_path), "--output", str(output_path)]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_docx_cli_fails_when_output_not_generated(tmp_path, monkeypatch):
    """個別CLIのdocxが、write_docxが何も書き出さない（例外を投げずに空振りする）場合に
    非ゼロ終了することを確認する。"""
    import src.cli as cli_module

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(cli_module, "write_docx", lambda path, document: None)
    output_path = tmp_path / "brushup.docx"
    monkeypatch.setattr(
        "sys.argv", ["cli", "docx", "--input", str(editable_path), "--output", str(output_path)]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_pdf_cli_fails_when_output_not_generated(tmp_path, monkeypatch):
    """個別CLIのpdfが、write_pdfが何も書き出さない場合に非ゼロ終了することを確認する。"""
    import src.cli as cli_module

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.setattr(cli_module, "write_pdf", lambda path, document: None)
    output_path = tmp_path / "brushup.pdf"
    monkeypatch.setattr(
        "sys.argv", ["cli", "pdf", "--input", str(editable_path), "--output", str(output_path)]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_canva_cli_succeeds_when_output_generated_normally(tmp_path, monkeypatch):
    """通常どおり出力が生成される場合は、canvaコマンドが正常終了することを確認する
    （成果物チェックが正常系を巻き添えにしないことの確認）。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_path = tmp_path / "canva_design.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "canva", "--input", str(editable_path), "--output", str(output_path)]
    )
    main()

    assert output_path.exists()
    assert output_path.stat().st_size > 0


# --- Phase 10.2追加修正: OCR必須モードのエラー表示重複整理 ---------------------------


def test_build_all_shows_single_consolidated_error_when_tesseract_missing(tmp_path, monkeypatch, capsys):
    """画像input + proofread + Tesseract未導入で、stderrに同じ意味の警告/エラーが
    重複表示されない（build-all側の集約エラー1ブロックのみ）ことを確認する。"""
    import src.cli as cli_module
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: not_ready_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1

    captured = capsys.readouterr()
    # build-all側の集約エラーは1回だけ出る。
    assert captured.err.count("mode=proofread requires OCR text") == 1
    # import_source()側の重複メッセージ（"OCR requires Tesseract"という個別の見出し）は出ない。
    assert "OCR requires Tesseract, but the 'tesseract' command was not found" not in captured.err
    # import_source()側の全ページ空警告も重複して出ない。
    assert "OCR produced no text for any page" not in captured.err


def test_build_all_allow_empty_ocr_shows_single_warning_not_import_source_warning(tmp_path, monkeypatch, capsys):
    """--allow-empty-ocr指定時も、import_source()側の個別警告ではなく、build-all側の
    集約された1つの警告だけが表示されることを確認する。"""
    import src.cli as cli_module
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: not_ready_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all", "--input", str(source_dir), "--mode", "proofread",
            "--output-dir", str(output_dir), "--allow-empty-ocr",
        ],
    )
    main()

    captured = capsys.readouterr()
    assert captured.err.count("OCR environment is degraded") == 1
    assert "OCR requires Tesseract, but the 'tesseract' command was not found" not in captured.err
    assert "OCR produced no text for any page" not in captured.err
    assert (output_dir / "editable" / "lesson_pages.json").exists()


def test_import_source_standalone_still_shows_warnings_when_tesseract_missing(tmp_path, monkeypatch, capsys):
    """単体のimport-sourceコマンドでは、従来どおりOCR関連の警告が表示されることを確認する
    （build-all向けの重複抑制がimport-source単体には影響しないことの回帰確認）。"""
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(import_source_module, "_try_ocr", lambda image_path, ocr_status: "")

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_path = tmp_path / "imported_pages.json"
    monkeypatch.setattr(
        "sys.argv", ["cli", "import-source", "--input", str(source_dir), "--output", str(output_path)]
    )
    main()

    captured = capsys.readouterr()
    assert "OCR requires Tesseract, but the 'tesseract' command was not found" in captured.err
    assert "OCR produced no text for any page" in captured.err


def test_build_all_failure_log_still_contains_ocr_details_despite_quiet_stderr(tmp_path, monkeypatch):
    """stderr表示は集約されていても、logs/には原因・警告・exit_codeが記録されることを確認する。"""
    import os

    import src.cli as cli_module
    import src.import_source as import_source_module

    not_ready_status = {
        "tesseract_available": False, "tesseract_path": None, "tesseract_on_path": False,
        "languages": [], "japanese_available": False, "english_available": False,
        "brew_available": False, "brew_path": None, "brew_on_path": False,
        "path_suggestions": [], "warnings": [], "errors": ["Tesseract command was not found."],
        "ocr_ready": False,
    }
    monkeypatch.setattr(import_source_module, "get_ocr_environment_status", lambda: not_ready_status)
    monkeypatch.setattr(cli_module, "get_ocr_environment_status", lambda: not_ready_status)

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit):
        main()

    from pathlib import Path

    log_files = list(Path(os.environ["AI_KYOUZAI_LOGS_DIR"]).glob("*_build-all.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "tesseract_available: False" in text
    assert "mode=proofread requires OCR text" in text
    assert "exit_code: 1" in text


# --- Phase 10.2追加修正: renderedディレクトリ空でも正常終了になる問題の修正 -------------


def test_build_all_fails_when_stale_rendered_files_exist_but_none_generated_this_run(tmp_path, monkeypatch):
    """output/rendered/に前回実行の古い画像ファイルが残っている状態で、今回の実行が
    画像を1枚も生成しなかった場合、非ゼロ終了になることを確認する（ディレクトリの
    非空判定だけでは検知できない、今回実行時点でのレンダリング失敗を検知するための回帰テスト）。"""
    import src.cli as cli_module

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"

    # 前回実行の古い画像ファイルをあらかじめ用意しておく。
    stale_dir = output_dir / "rendered"
    stale_dir.mkdir(parents=True, exist_ok=True)
    (stale_dir / "page_001.png").write_bytes(b"stale image from a previous run")

    # 今回の実行ではrender_document_images()が1枚も生成しなかった状況を再現する。
    monkeypatch.setattr(cli_module, "render_document_images", lambda *args, **kwargs: [])

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_build_all_fails_when_rendered_image_is_zero_bytes(tmp_path, monkeypatch):
    """render_document_images()が返したパスの画像がサイズ0の場合、非ゼロ終了になることを確認する。"""
    import src.cli as cli_module

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=1)
    output_dir = tmp_path / "output"

    def _fake_render(document, out_dir, rendered_dir, font_path=None):
        rendered_dir = Path(rendered_dir)
        rendered_dir.mkdir(parents=True, exist_ok=True)
        zero_path = rendered_dir / "page_001.png"
        zero_path.write_bytes(b"")
        return [zero_path]

    monkeypatch.setattr(cli_module, "render_document_images", _fake_render)

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_build_all_fails_when_rendered_count_does_not_match_page_count(tmp_path, monkeypatch):
    """render_document_images()が返した画像数がページ数と一致しない場合、非ゼロ終了になることを
    確認する。"""
    import src.cli as cli_module

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=3)
    output_dir = tmp_path / "output"

    def _fake_render(document, out_dir, rendered_dir, font_path=None):
        rendered_dir = Path(rendered_dir)
        rendered_dir.mkdir(parents=True, exist_ok=True)
        only_one = rendered_dir / "page_001.png"
        only_one.write_bytes(b"only one image, but there are 3 pages")
        return [only_one]

    monkeypatch.setattr(cli_module, "render_document_images", _fake_render)

    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_build_all_image_input_same_format_generates_rendered_images_successfully(tmp_path, monkeypatch):
    """画像input + --output-format same（既定）で、rendered画像が実際に生成されることを
    確認する（正常系が壊れていないことの確認）。"""
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=3)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    rendered_files = sorted((output_dir / "rendered").glob("*.png"))
    assert len(rendered_files) == 3
    for f in rendered_files:
        assert f.stat().st_size > 0


def test_build_all_output_format_image_generates_rendered_images_successfully(tmp_path, monkeypatch):
    """--output-format imageで、rendered画像が1枚以上生成されることを確認する。"""
    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all", "--input", str(source_dir), "--mode", "proofread",
            "--output-dir", str(output_dir), "--output-format", "image",
        ],
    )
    main()

    rendered_files = sorted((output_dir / "rendered").glob("*.png"))
    assert len(rendered_files) == 2


def test_build_all_output_format_all_validates_rendered_images(tmp_path, monkeypatch):
    """--output-format allでも、rendered画像が検証対象になり、正しく生成されることを確認する。"""
    import src.cli as cli_module

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=2)
    output_dir = tmp_path / "output"

    # allの中でrenderedだけ空振りするケースをシミュレートし、検証されることを確認する。
    monkeypatch.setattr(cli_module, "render_document_images", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "build-all", "--input", str(source_dir), "--mode", "proofread",
            "--output-dir", str(output_dir), "--output-format", "all",
        ],
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


def test_build_all_logs_rendered_files_count_not_just_directory(tmp_path, monkeypatch):
    """ログのgenerated_files/RENDERED_IMAGESセクションに、ディレクトリ名だけでなく
    実ファイルパスまたは生成枚数が記録されることを確認する。"""
    import os

    source_dir = tmp_path / "source"
    _make_source_images(source_dir, count=3)
    output_dir = tmp_path / "output"
    monkeypatch.setattr(
        "sys.argv",
        ["cli", "build-all", "--input", str(source_dir), "--mode", "proofread", "--output-dir", str(output_dir)],
    )
    main()

    log_files = list(Path(os.environ["AI_KYOUZAI_LOGS_DIR"]).glob("*_build-all.log"))
    assert len(log_files) == 1
    text = log_files[0].read_text(encoding="utf-8")
    assert "rendered_files_count: 3" in text
    assert "page_001.png" in text


def test_regenerate_fails_when_rendered_images_not_actually_generated(tmp_path, monkeypatch):
    """regenerateでも、rendered画像の検証がstale-file問題を回避することを確認する。"""
    import src.cli as cli_module

    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_dir = editable_path.parent.parent
    stale_dir = output_dir / "rendered"
    stale_dir.mkdir(parents=True, exist_ok=True)
    (stale_dir / "page_001.png").write_bytes(b"stale image from a previous run")

    monkeypatch.setattr(cli_module, "render_document_images", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        "sys.argv", ["cli", "regenerate", "--input", str(editable_path), "--output-format", "image"]
    )
    with pytest.raises(SystemExit) as exc_info:
        main()
    assert exc_info.value.code == 1


# --- LLM手作業投入用中間ファイル生成（llm-handoff） ---------------------------------


def test_llm_handoff_cli_generates_markdown_from_editable_lesson_pages(tmp_path, monkeypatch):
    """build-all/lesson-pagesで作ったeditable/lesson_pages.jsonから、llm-handoffコマンドで
    llm_handoff.mdが生成されることを確認する（CLIからの実行確認）。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_path = tmp_path / "output" / "llm_handoff.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "llm-handoff", "--input", str(editable_path), "--output", str(output_path)]
    )
    main()

    assert output_path.exists()
    assert output_path.stat().st_size > 0
    text = output_path.read_text(encoding="utf-8")
    assert "ブラッシュアップであって、作り直しではない" in text
    assert "### Page 1" in text


def test_llm_handoff_cli_default_output_path(tmp_path, monkeypatch):
    """--outputを省略した場合、output/llm_handoff.md（カレントディレクトリ基準）に
    生成されることを確認する。"""
    editable_path = tmp_path / "editable_lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("sys.argv", ["cli", "llm-handoff", "--input", str(editable_path)])
    main()

    assert (tmp_path / "output" / "llm_handoff.md").exists()


# --- LLM回答後の手作業フロー：採用判断シート（edit-plan-template） -----------------------


def test_edit_plan_template_cli_generates_markdown_from_editable_lesson_pages(tmp_path, monkeypatch):
    """build-all/lesson-pagesで作ったeditable/lesson_pages.jsonから、edit-plan-templateコマンドで
    edit_plan_template.mdが生成されることを確認する（CLIからの実行確認）。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_path = tmp_path / "output" / "edit_plan_template.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "edit-plan-template", "--input", str(editable_path), "--output", str(output_path)]
    )
    main()

    assert output_path.exists()
    assert output_path.stat().st_size > 0
    text = output_path.read_text(encoding="utf-8")
    assert "採用判断" in text
    assert "### Page 1" in text


def test_edit_plan_template_cli_default_output_path(tmp_path, monkeypatch):
    """--outputを省略した場合、output/edit_plan_template.md（カレントディレクトリ基準）に
    生成されることを確認する。"""
    editable_path = tmp_path / "editable_lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("sys.argv", ["cli", "edit-plan-template", "--input", str(editable_path)])
    main()

    assert (tmp_path / "output" / "edit_plan_template.md").exists()


def test_llm_handoff_md_includes_guidance_to_use_edit_plan_template(tmp_path, monkeypatch):
    """llm_handoff.mdの注意事項に、LLM回答後はedit-plan-templateで採用判断を整理する案内が
    含まれることを確認する。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_path = tmp_path / "output" / "llm_handoff.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "llm-handoff", "--input", str(editable_path), "--output", str(output_path)]
    )
    main()

    text = output_path.read_text(encoding="utf-8")
    assert "edit-plan-template" in text
    assert "edit_plan_template.md" in text


# --- OCR品質チェック＋補正候補データ生成（ocr-check） ------------------------------------


def test_ocr_check_cli_generates_report_and_candidates_json(tmp_path, monkeypatch):
    """build-all/lesson-pagesで作ったeditable/lesson_pages.jsonから、ocr-checkコマンドで
    Markdownレポートと補正候補JSONの両方が生成されることを確認する（CLIからの実行確認）。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    report_path = tmp_path / "output" / "ocr_check_report.md"
    candidates_path = tmp_path / "output" / "ocr_correction_candidates.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "ocr-check", "--input", str(editable_path),
            "--output", str(report_path), "--candidates-output", str(candidates_path),
        ],
    )
    main()

    assert report_path.exists()
    assert report_path.stat().st_size > 0
    assert candidates_path.exists()
    assert candidates_path.stat().st_size > 0

    import json
    data = json.loads(candidates_path.read_text(encoding="utf-8"))
    assert "candidates" in data
    assert "summary" in data


def test_ocr_check_cli_custom_candidates_output_path(tmp_path, monkeypatch):
    """--candidates-outputで補正候補JSONの出力先を変更できることを確認する。"""
    editable_path = tmp_path / "editable_lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    custom_candidates_path = tmp_path / "custom_candidates.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "ocr-check", "--input", str(editable_path),
            "--candidates-output", str(custom_candidates_path),
        ],
    )
    main()

    assert custom_candidates_path.exists()


def test_ocr_check_cli_default_output_paths(tmp_path, monkeypatch):
    """--output/--candidates-outputを省略した場合、既定パスに生成されることを確認する。"""
    editable_path = tmp_path / "editable_lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    monkeypatch.chdir(tmp_path)
    monkeypatch.setattr("sys.argv", ["cli", "ocr-check", "--input", str(editable_path)])
    main()

    assert (tmp_path / "output" / "ocr_check_report.md").exists()
    assert (tmp_path / "output" / "ocr_correction_candidates.json").exists()


def test_llm_handoff_md_includes_ocr_check_guidance(tmp_path, monkeypatch):
    """llm_handoff.mdの注意事項に、LLM投入前にocr-checkを確認する案内が含まれることを確認する。"""
    editable_path = tmp_path / "output" / "editable" / "lesson_pages.json"
    monkeypatch.setattr(
        "sys.argv",
        [
            "cli", "lesson-pages", "--mode", "generate",
            "--requirements", "examples/requirements_ai_instagram.json",
            "--output", str(editable_path),
        ],
    )
    main()

    output_path = tmp_path / "output" / "llm_handoff.md"
    monkeypatch.setattr(
        "sys.argv", ["cli", "llm-handoff", "--input", str(editable_path), "--output", str(output_path)]
    )
    main()

    text = output_path.read_text(encoding="utf-8")
    assert "ocr-check" in text
